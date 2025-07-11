import os
os.environ["OLLAMA_NO_CUDA"] = "1"

import subprocess
import socket
import tempfile
import time
import xml.etree.ElementTree as ET
from pathlib import Path

import streamlit as st
import docx
from pptx import Presentation
import pandas as pd
import ollama
import chromadb
from ultralytics import YOLO
import matplotlib.pyplot as plt

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyMuPDFLoader
from chromadb.utils.embedding_functions.ollama_embedding_function import OllamaEmbeddingFunction

# ----------- Start Ollama Automatically ----------- #
def is_ollama_running():
    s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    try:
        s.connect(("localhost", 11434))
        s.shutdown(socket.SHUT_RDWR)
        return True
    except:
        return False
    finally:
        s.close()

if not is_ollama_running():
    try:
        subprocess.Popen(["ollama", "serve"], creationflags=subprocess.CREATE_NO_WINDOW)

        time.sleep(2)
    except Exception as e:
        st.error(f"❌ Failed to start Ollama: {e}")
        st.stop()

# ------------------ UI + Styling ------------------ #
st.set_page_config(page_title="InsightIQ", layout="centered")
st.markdown("""
    <style>
    html, body, [data-testid="stApp"] {
        background-color: #121212;
        color: #f1f1f1;
        font-family: 'Segoe UI', sans-serif;
    }
    .title-container { text-align: center; padding: 1rem 0; }
    .title-container h1 {
        font-size: 2.5rem;
        font-weight: bold;
        color: #8a7ff0;
        margin-bottom: 0;
    }
    .stTextInput>div>div>input {
        background-color: #1e1e1e;
        color: #fff;
    }
    .stButton button {
        background-color: #8a7ff0 !important;
        color: white;
        border-radius: 10px;
    }
    </style>
    <div class='title-container'><h1>🤖InsightIQ</h1></div>
""", unsafe_allow_html=True)

# ------------------ Session State ------------------ #
if "collection" not in st.session_state:
    st.session_state.collection = None
if "all_docs" not in st.session_state:
    st.session_state.all_docs = []

# ------------------ File Processors ------------------ #
def create_document_with_path(text, file_path):
    return Document(page_content=text, metadata={"source": str(file_path)})

def process_pdf(file, file_path):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file.read())
    loader = PyMuPDFLoader(tmp.name)
    docs = loader.load()
    os.unlink(tmp.name)
    for doc in docs:
        doc.metadata["source"] = str(file_path)
    return docs

def process_docx(file, file_path):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(file.read())
        doc = docx.Document(tmp.name)
        os.unlink(tmp.name)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs[:100])
        return [Document(page_content=text, metadata={"source": str(file_path)})]
    except:
        return []

def process_pptx(file, file_path):
    import io
    from PIL import Image

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(file.read())
    prs = Presentation(tmp.name)
    os.unlink(tmp.name)

    texts = []
    detected_objects = []

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

            # Extract images
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                img_bytes = image.blob
                with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as img_tmp:
                    img_tmp.write(img_bytes)
                    img_path = img_tmp.name

                # Run YOLO detection
                try:
                    model = YOLO("yolov8m.pt")
                    results = model(img_path, conf=0.5)[0]
                    labels = [model.names[int(cls)] for cls in results.boxes.cls]
                    obj_count = {label: labels.count(label) for label in set(labels)}
                    detected_text = f"[Slide {i+1} - Detected Objects]\n" + "\n".join(f"{label}: {count}" for label, count in obj_count.items())
                    detected_objects.append(detected_text)
                except Exception as e:
                    detected_objects.append(f"[Slide {i+1}] ⚠️ Error detecting image: {e}")
                finally:
                    os.unlink(img_path)

    full_text = "\n".join(texts + detected_objects)
    return [Document(page_content=full_text, metadata={"source": str(file_path)})]



def process_xml(file, file_path):
    try:
        tree = ET.ElementTree(ET.fromstring(file.read()))
        root = tree.getroot()
        def extract(elem, level=0):
            indent = "  " * level
            lines = [f"{indent}{elem.tag}: {elem.text.strip() if elem.text else ''}"]
            for child in elem:
                lines.extend(extract(child, level + 1))
            return lines
        content = extract(root)
        return [create_document_with_path("[XML CONTENT]\n" + "\n".join(content), file_path)]
    except:
        return []

def process_txt(file, file_path):
    try:
        content = file.read().decode("utf-8", errors="ignore")
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        return [create_document_with_path("\n".join(lines), file_path)]
    except:
        return []

def process_csv(file, file_path):
    try:
        df = pd.read_csv(file)
        st.session_state.all_docs = [df]  # Only store the latest for chart
        return [create_document_with_path(df.to_string(index=False), file_path)]
    except:
        return []

def process_xlsx(file, file_path):
    try:
        df = pd.read_excel(file, engine="openpyxl")
        st.session_state.all_docs = [df]  # Only store the latest for chart
        return [create_document_with_path(df.to_string(index=False), file_path)]
    except:
        return []

def process_image(file, file_path):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
            tmp.write(file.read())
        path = tmp.name
        model = YOLO("yolov8m.pt")
        results = model(path, conf=0.5)[0]
        labels = [model.names[int(cls)] for cls in results.boxes.cls]
        os.unlink(path)
        obj_count = {label: labels.count(label) for label in set(labels)}
        obj_text = "\n".join(f"{label}: {count}" for label, count in obj_count.items())
        safe_text = obj_text.encode("utf-16", "surrogatepass").decode("utf-16", "ignore")
        full_text = "[Detected Objects]\n" + safe_text if safe_text.strip() else ""

        return [Document(page_content=full_text, metadata={"source": str(file_path)})]
    except:
        return []

# ------------------ Embedding & Vector Store ------------------ #
def get_vectorstore(documents):
    client = chromadb.Client()
    ollama_embed = OllamaEmbeddingFunction(
        model_name="nomic-embed-text",
        url="http://localhost:11434/api/embeddings"
    )
    collection = client.get_or_create_collection("insightiq_docs", embedding_function=ollama_embed)
    existing_ids = collection.get()["ids"]
    if existing_ids:
        collection.delete(ids=existing_ids)
    for i, doc in enumerate(documents):
        metadata = doc.metadata if doc.metadata else {"source": "unknown"}
        collection.add(documents=[doc.page_content], metadatas=[metadata], ids=[str(i)])
    return collection

# ------------------ QA System ------------------ #
def ask_llm(user_input, collection):
    with st.spinner("🤔 Thinking..."):
        results = collection.query(query_texts=[user_input], n_results=10)
        if not results["documents"] or not results["documents"][0]:
            return "❌ No relevant content found."
        context_chunks = results["documents"][0]
        sources = results["metadatas"][0]
        context = ""
        for i, chunk in enumerate(context_chunks):
            source = sources[i].get("source", "unknown")
            context += f"[From: {source}]\n{chunk}\n\n"
        prompt = f"""
You are a helpful assistant. Use the document excerpts below to answer the user's question clearly and uniquely.

Context:
{context}

Question: {user_input}
"""
        response = ollama.chat(model="llama3", messages=[{"role": "user", "content": prompt}])
        return response["message"]["content"]

# ------------------ Upload UI ------------------ #
MAX_FILES = 10
with st.expander("📁 Upload supported files", expanded=True):
    uploaded_files = st.file_uploader(
        "Upload your files below:",
        type=["pdf", "docx", "pptx", "xml", "txt", "jpg", "jpeg", "png", "csv", "xlsx"],
        accept_multiple_files=True
    )
    if uploaded_files:
        if len(uploaded_files) > MAX_FILES:
            st.warning(f"⚠️ Maximum {MAX_FILES} files allowed.")
        else:
            docs = []
            for file in uploaded_files:
                ext = file.name.split(".")[-1].lower()
                file_path = os.path.abspath(file.name)
                st.write(f"📄 File uploaded from: `{file_path}`")
                if ext == "pdf": docs.extend(process_pdf(file, file_path))
                elif ext == "docx": docs.extend(process_docx(file, file_path))
                elif ext == "pptx": docs.extend(process_pptx(file, file_path))
                elif ext == "xml": docs.extend(process_xml(file, file_path))
                elif ext == "txt": docs.extend(process_txt(file, file_path))
                elif ext == "csv": docs.extend(process_csv(file, file_path))
                elif ext == "xlsx": docs.extend(process_xlsx(file, file_path))
                elif ext in ["jpg", "jpeg", "png"]: docs.extend(process_image(file, file_path))
            if docs:
                chunks = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=80).split_documents(docs)
                st.session_state.collection = get_vectorstore(chunks)
                st.success("✅ Files processed and embedded successfully!")

# ------------------ Visualization ------------------ #
def try_visualize(query):
    if not st.session_state.all_docs:
        return

    df = st.session_state.all_docs[-1]
    if "country" in df.columns.str.lower().tolist():
        country_col = [col for col in df.columns if col.lower() == "country"][0]
        country_counts = df[country_col].value_counts()

        if "bar chart" in query.lower():
            st.subheader("🌍 Country Distribution")
            st.bar_chart(country_counts)

        elif "pie chart" in query.lower():
            st.subheader("🌍 Country Distribution")
            fig, ax = plt.subplots()
            ax.pie(country_counts.head(10), labels=country_counts.head(10).index, autopct='%1.1f%%')
            ax.axis('equal')
            st.pyplot(fig)

# ------------------ Chat Input ------------------ #
# ------------------ Display Chat History ------------------ #
for entry in st.session_state.get("chat_history", []):
    st.chat_message("user").write(entry["user"])
    st.chat_message("assistant").write(entry["bot"])

# ------------------ Chat Input ------------------ #
user_input = st.chat_input("Ask a question from the uploaded documents or images...")
if user_input:
    if not st.session_state.collection:
        st.error("📁 Please upload and process documents first.")
    else:
        st.chat_message("user").write(user_input)
        start = time.time()
        answer = ask_llm(user_input, st.session_state.collection)
        st.chat_message("assistant").write(answer)
        st.caption(f"🕒 Response time: {time.time() - start:.2f} seconds")

        # Save to chat history
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []
        st.session_state.chat_history.append({"user": user_input, "bot": answer})

        try_visualize(user_input)
